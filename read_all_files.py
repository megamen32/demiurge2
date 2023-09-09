import PyPDF2

from docx import Document
import openpyxl
from pptx import Presentation

def read_file(file_path):
    file_extension = file_path.split('.')[-1]
    content = ""

    if file_extension == 'pdf':
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                content += page.extract_text()


    elif file_extension in ['doc', 'docx']:
        doc = Document(file_path)
        for para in doc.paragraphs:
            content += para.text + '\n'

    elif file_extension in ['xls', 'xlsx']:
        wb = openpyxl.load_workbook(file_path)
        sheet = wb.active
        for row in sheet.iter_rows():
            for cell in row:
                content += str(cell.value) + '\t'
            content += '\n'

    elif file_extension in ['ppt', 'pptx']:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            content += run.text + ' '
    else:
        with open(file_path, 'r', encoding='utf-8') as f:
            content= f.read()

    return content
